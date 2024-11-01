import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import argparse

class Variables:
    input_file = 'input.xlsx'
    output_file = 'quiz_presentation.pptx' 
    title_slide_text = "Кто самый не тупой"  # Main slide header
    back_button_text = "Назад" 
    CategoryColumn = 'Категория'
    question_slide_layout = 1  # Questions Slide Index
    category_slide_layout = 5  # Categories Slide Index
    QuestionsSlideNumber = 3 
    MainSlideName = 'MainTable'
    TemplatePresentationPath = "Template\\Template.pptx"
    TableStyleGUID =  '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'  #NoStyleNoGrid


def ConvertRelativePathToAbsolute(path):
    return os.path.abspath(path) if not os.path.isabs(path) else path

def read_excel_table(file_path, sheet_name=None):
    """Считывает данные из указанного листа Excel файла."""
    try:
        if sheet_name is None:
            df = pd.read_excel(file_path)
        else:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
        
        # Проверка на наличие данных в DataFrame
        if not isinstance(df, pd.DataFrame):
            raise ValueError("Ошибка: Полученные данные не являются DataFrame.")
        
        if df.empty:
            raise ValueError("Ошибка: Таблица пустая. Проверьте содержимое файла Excel.")
        
        return df
    except ValueError as e:
        raise ValueError(f"Ошибка при чтении Excel файла: {e}")
    except Exception as e:
        raise Exception(f"Общая ошибка: {e}")

def CreateQuizPresentation(df, output_file=Variables.output_file):
    Slides = []

    #Open .pptx Template
    prs = Presentation(ConvertRelativePathToAbsolute(Variables.TemplatePresentationPath)) 

    #Set slide params               
    slide_width = prs.slide_width 
    slide_height = prs.slide_height
    left = Inches(0.5)
    top = Inches(0.5)
    width = slide_width - Inches(1)
    height = slide_height - Inches(1)

    # Create main slide
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = Variables.title_slide_text

    # Create Categories Slides

    for CategoryIndex in range(len(df[Variables.CategoryColumn])):
        question_slide = prs.slides.add_slide(prs.slide_layouts[Variables.question_slide_layout])
        question_title = question_slide.shapes.title
        question_content = question_slide.shapes.placeholders[1]
        question_title.text = f"{str(df[Variables.CategoryColumn][CategoryIndex])}"

    # Create Table Slide
    category_slide_layout = prs.slide_layouts[Variables.category_slide_layout]
    category_slide = prs.slides.add_slide(category_slide_layout)
    question_title = category_slide.shapes.title
    question_title.text = f"{str(Variables.MainSlideName)}"
    question_title.text_frame.paragraphs[0].font.size = Pt(1)

    # Create Table
    shape = category_slide.shapes.add_table(len(df), len(df.columns), left, top, width, height)
    table = shape.table
    tbl =  shape._element.graphic.graphicData.tbl
    style_id = Variables.TableStyleGUID
    tbl[0][-1].text = style_id

  
    #Get Start Slide For Questions

    StartSlide = Variables.QuestionsSlideNumber + len(df[Variables.CategoryColumn]) 

    #Set Categories in table

    for CategoryIndex in range(len(df[Variables.CategoryColumn])): #Set Categories in table 
        table.cell(CategoryIndex, 0).text = str(df[Variables.CategoryColumn][CategoryIndex])
    
    #Get Slides Param
    for index, row in df.iterrows():
        for column in [200, 400, 600, 800, 1000]:  #Columns
            match column:
                case 200:
                    XCoordinate = 0
                case 400:
                    XCoordinate = 1
                case 600:
                    XCoordinate = 2
                case 800:
                    XCoordinate = 3
                case 1000:
                    XCoordinate = 4
           
            SlideTemplate = {
                'Header' : df[Variables.CategoryColumn][index] ,
                'Text': row[column],
                'SlideNumber': StartSlide,
                'X': XCoordinate+1,
                'Y': index,
                'TextCell': column,
                'Link': '',
            }
            StartSlide += 1
            Slides.append(SlideTemplate)

    #Creating slides with questions
    for slide in Slides: 
        question_slide = prs.slides.add_slide(prs.slide_layouts[Variables.question_slide_layout])
        question_title = question_slide.shapes.title
        question_content = question_slide.shapes.placeholders[1]
        question_title.text = f"{slide['Header']} - {slide['TextCell']}"
        question_content.text = f"{slide['Text']}"

        #Creating "return to main slide hyperlink"
        txBox = question_slide.shapes.add_textbox(Inches(7), Inches(6.5), Inches(2), Inches(1))
        tf = txBox.text_frame
        link = tf.paragraphs[0].add_run()
        link.text = Variables.back_button_text
        link.font.size = Pt(11)
        link.hyperlink.address = f"#{Variables.MainSlideName}"
    
    #Add hyperlinks to table
    for slide in Slides: 
        cell = table.cell(slide['Y'], slide['X'])
        cell_link = cell.text_frame.paragraphs[0].add_run()
        cell_link.text = str(slide['TextCell'])
        hlink = cell_link.hyperlink
        hlink.address = f"#{slide['Header']} - {slide['TextCell']}"

    #Save Presentation
    try:
        prs.save(output_file)
        print(f"Presentation saved: {output_file}")
    except PermissionError:
         print(f"No access to file: {output_file}")

def GetArgs():
    parser = argparse.ArgumentParser(description='Создание презентации викторины из Excel файла.')
    parser.add_argument('--sheet', type=str, default=None, 
                        help='Имя листа Excel (по умолчанию используется первый лист).')
    return parser.parse_args()

def main():
    args = GetArgs()
    CreateQuizPresentation(read_excel_table(ConvertRelativePathToAbsolute(Variables.input_file), sheet_name=args.sheet))

if __name__ == "__main__":
    main()
